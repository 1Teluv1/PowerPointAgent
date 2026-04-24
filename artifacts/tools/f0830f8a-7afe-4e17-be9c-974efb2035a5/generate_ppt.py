import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


OUTPUT_FILE = "perfume_brand_introduction.pptx"

SLIDE_W = 13.333333
SLIDE_H = 7.5

COLORS = {
    "background": "0E1726",
    "surface": "F7F1E8",
    "accent": "C9A45C",
    "secondary": "8A6F48",
    "text": "1A1A1A",
    "muted": "7A746C",
    "white": "FFFFFF",
    "soft_navy": "172338",
    "line": "D8C69B",
}

FONTS = {
    "title": "Aptos Display",
    "body": "Aptos",
    "fallback": "Calibri",
}

ASSETS = {
    "brand_essence_triads": "assets/brand_essence_triads.svg",
    "market_opportunity_kpis": "assets/market_opportunity_kpis.svg",
    "customer_scent_journey": "assets/customer_scent_journey.svg",
    "product_lineup_comparison": "assets/product_lineup_comparison.svg",
    "scent_note_pyramid": "assets/scent_note_pyramid.svg",
    "brand_positioning_map": "assets/brand_positioning_map.svg",
    "launch_roadmap": "assets/launch_roadmap.svg",
}


def rgb(hex_color):
    value = hex_color.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def inches(value):
    return Inches(value)


def set_slide_background(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size=18,
    color=None,
    font_name=None,
    bold=False,
    italic=False,
    align=None,
    valign=None,
    line_spacing=1.0
):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    if valign is not None:
        frame.vertical_anchor = valign

    paragraph = frame.paragraphs[0]
    paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name or FONTS["body"]
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color or COLORS["text"])
    return box


def add_shape(slide, shape_type, x, y, w, h, fill_color=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, inches(x), inches(y), inches(w), inches(h))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1, y1, x2, y2, color=None, width=1):
    line = slide.shapes.add_connector(1, inches(x1), inches(y1), inches(x2), inches(y2))
    line.line.color.rgb = rgb(color or COLORS["accent"])
    line.line.width = Pt(width)
    return line


def add_divider(slide, x, y, w, color=None, width=1):
    return add_line(slide, x, y, x + w, y, color=color or COLORS["accent"], width=width)


def add_card(slide, x, y, w, h, fill_color=None, line_color=None):
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill_color=fill_color or COLORS["surface"],
        line_color=line_color or COLORS["line"],
        line_width=1
    )


def add_missing_asset_placeholder(slide, asset_id, x, y, w, h):
    add_card(slide, x, y, w, h, fill_color="132033", line_color=COLORS["accent"])
    add_text(
        slide,
        "Missing visual asset",
        x + 0.25,
        y + h / 2 - 0.28,
        w - 0.5,
        0.3,
        font_size=13,
        color=COLORS["surface"],
        font_name=FONTS["body"],
        bold=True,
        align=PP_ALIGN.CENTER
    )
    add_text(
        slide,
        asset_id,
        x + 0.25,
        y + h / 2 + 0.08,
        w - 0.5,
        0.3,
        font_size=10,
        color=COLORS["accent"],
        font_name=FONTS["body"],
        align=PP_ALIGN.CENTER
    )


def add_image_or_placeholder(slide, asset_id, path, x, y, w, h):
    if os.path.exists(path):
        try:
            slide.shapes.add_picture(path, inches(x), inches(y), width=inches(w), height=inches(h))
        except Exception:
            add_missing_asset_placeholder(slide, asset_id, x, y, w, h)
    else:
        add_missing_asset_placeholder(slide, asset_id, x, y, w, h)


def add_small_label(slide, text, x, y, w, color=None):
    add_text(
        slide,
        text,
        x,
        y,
        w,
        0.24,
        font_size=9,
        color=color or COLORS["accent"],
        font_name=FONTS["body"],
        bold=True,
        align=PP_ALIGN.LEFT
    )


def add_corner_mark(slide):
    add_line(slide, 0.55, 0.55, 1.25, 0.55, color=COLORS["accent"], width=1)
    add_line(slide, 0.55, 0.55, 0.55, 1.25, color=COLORS["accent"], width=1)
    add_line(slide, 12.1, 6.95, 12.8, 6.95, color=COLORS["accent"], width=1)
    add_line(slide, 12.8, 6.25, 12.8, 6.95, color=COLORS["accent"], width=1)


def add_footer(slide, text="Scent Atelier"):
    add_text(
        slide,
        text,
        0.65,
        7.03,
        2.2,
        0.24,
        font_size=8,
        color=COLORS["secondary"],
        font_name=FONTS["body"]
    )


def add_perfume_bottle_icon(slide, x, y, scale=1.0, color=None):
    fill = color or COLORS["accent"]
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.18 * scale, y, 0.24 * scale, 0.18 * scale, fill_color=fill, line_color=fill)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.12 * scale, y + 0.18 * scale, 0.36 * scale, 0.16 * scale, fill_color=COLORS["surface"], line_color=fill)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.34 * scale, 0.6 * scale, 0.78 * scale, fill_color=None, line_color=fill)
    add_line(slide, x + 0.12 * scale, y + 0.74 * scale, x + 0.48 * scale, y + 0.74 * scale, color=fill, width=1)


def add_section_title(slide, title, subtitle, dark=True):
    title_color = COLORS["surface"] if dark else COLORS["background"]
    subtitle_color = COLORS["accent"] if dark else COLORS["secondary"]
    add_small_label(slide, "BRAND INTRODUCTION", 0.72, 0.55, 2.2, color=subtitle_color)
    add_text(
        slide,
        title,
        0.72,
        0.88,
        5.3,
        0.58,
        font_size=30,
        color=title_color,
        font_name=FONTS["title"],
        bold=True
    )
    add_text(
        slide,
        subtitle,
        0.74,
        1.48,
        5.2,
        0.35,
        font_size=13,
        color=subtitle_color,
        font_name=FONTS["body"]
    )
    add_divider(slide, 0.72, 1.92, 1.1, color=subtitle_color, width=1)


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])

    add_shape(slide, MSO_SHAPE.OVAL, 8.2, -1.1, 5.8, 5.8, fill_color="172338", line_color=None)
    add_shape(slide, MSO_SHAPE.OVAL, -1.2, 4.5, 4.1, 4.1, fill_color="172338", line_color=None)
    add_corner_mark(slide)

    add_text(
        slide,
        "Scent Atelier",
        2.0,
        2.55,
        9.33,
        0.78,
        font_size=47,
        color=COLORS["surface"],
        font_name=FONTS["title"],
        bold=True,
        align=PP_ALIGN.CENTER
    )
    add_text(
        slide,
        "A Premium Perfume Brand Introduction",
        3.1,
        3.35,
        7.1,
        0.35,
        font_size=16,
        color=COLORS["accent"],
        font_name=FONTS["body"],
        align=PP_ALIGN.CENTER
    )
    add_divider(slide, 5.62, 3.95, 2.1, color=COLORS["accent"], width=1)

    add_text(
        slide,
        "향으로 기억되는 순간을 디자인하는 프리미엄 니치 향수 브랜드",
        2.8,
        6.38,
        7.75,
        0.35,
        font_size=13,
        color=COLORS["surface"],
        font_name="Malgun Gothic",
        align=PP_ALIGN.CENTER
    )
    add_perfume_bottle_icon(slide, 6.37, 4.48, scale=1.0, color=COLORS["accent"])


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])
    add_section_title(slide, "Brand Essence", "향으로 완성하는 개인의 분위기", dark=True)

    add_text(
        slide,
        "Scent Atelier는 단순한 향이 아니라, 개인의 기억과 분위기를 정교하게 설계하는 브랜드입니다.",
        0.74,
        2.35,
        4.25,
        0.92,
        font_size=20,
        color=COLORS["surface"],
        font_name="Malgun Gothic",
        bold=True,
        line_spacing=1.08
    )
    add_text(
        slide,
        "브랜드는 감각적 기억, 절제된 우아함, 오래 남는 잔향이라는 세 가지 가치를 중심으로 전개됩니다.",
        0.76,
        3.58,
        3.95,
        1.05,
        font_size=13,
        color="D9D0C3",
        font_name="Malgun Gothic",
        line_spacing=1.08
    )

    values = [
        ("01", "Sensory Memory", "향으로 떠오르는 장면"),
        ("02", "Quiet Luxury", "과시보다 절제된 고급감"),
        ("03", "Lasting Impression", "시간이 지나도 남는 여운"),
    ]
    for i, item in enumerate(values):
        y = 4.95 + i * 0.52
        add_text(slide, item[0], 0.78, y, 0.35, 0.24, font_size=9, color=COLORS["accent"], font_name=FONTS["body"], bold=True)
        add_text(slide, item[1], 1.18, y - 0.02, 1.9, 0.26, font_size=11, color=COLORS["surface"], font_name=FONTS["body"], bold=True)
        add_text(slide, item[2], 3.05, y - 0.02, 2.0, 0.26, font_size=10, color="BFB6A8", font_name="Malgun Gothic")

    add_image_or_placeholder(slide, "brand_essence_triads", ASSETS["brand_essence_triads"], 5.55, 1.25, 7.05, 5.85)
    add_footer(slide)


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["surface"])

    add_section_title(slide, "Market Opportunity", "니치 향수 시장의 성장과 기회", dark=False)
    add_text(
        slide,
        "개성 표현, 프리미엄 소비, 선물 수요 증가가 니치 향수 브랜드의 진입 기회를 만듭니다.",
        6.15,
        0.84,
        5.85,
        0.68,
        font_size=16,
        color=COLORS["text"],
        font_name="Malgun Gothic",
        bold=True,
        line_spacing=1.08
    )

    add_image_or_placeholder(slide, "market_opportunity_kpis", ASSETS["market_opportunity_kpis"], 0.82, 2.05, 11.7, 4.25)

    add_card(slide, 2.1, 6.52, 9.05, 0.46, fill_color="FFFFFF", line_color="E5D8BA")
    add_text(
        slide,
        "Insight: 향수는 기능재를 넘어 취향·선물·자기표현을 결합한 프리미엄 라이프스타일 카테고리로 확장되고 있습니다.",
        2.35,
        6.64,
        8.55,
        0.22,
        font_size=9.5,
        color=COLORS["secondary"],
        font_name="Malgun Gothic",
        align=PP_ALIGN.CENTER
    )


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])
    add_section_title(slide, "Target Customer", "감각적 취향을 소비하는 25–39세 고객", dark=True)

    add_card(slide, 0.8, 2.28, 4.45, 4.25, fill_color=COLORS["surface"], line_color=COLORS["accent"])
    add_text(slide, "Core Persona", 1.08, 2.55, 1.8, 0.28, font_size=11, color=COLORS["secondary"], font_name=FONTS["body"], bold=True)
    add_text(slide, "Urban Taste Seeker", 1.08, 2.95, 3.45, 0.48, font_size=24, color=COLORS["background"], font_name=FONTS["title"], bold=True)
    add_text(
        slide,
        "도시적 라이프스타일과 자기표현 욕구를 가진 25–39세 고객",
        1.1,
        3.56,
        3.45,
        0.56,
        font_size=12,
        color=COLORS["text"],
        font_name="Malgun Gothic",
        line_spacing=1.05
    )

    persona_items = [
        ("Lifestyle", "전시·카페·호텔·편집숍 경험 선호"),
        ("Motivation", "나만의 분위기를 만드는 향 탐색"),
        ("Occasion", "데일리 사용과 프리미엄 선물"),
    ]
    for i, item in enumerate(persona_items):
        y = 4.45 + i * 0.52
        add_text(slide, item[0], 1.1, y, 1.15, 0.22, font_size=9.5, color=COLORS["accent"], font_name=FONTS["body"], bold=True)
        add_text(slide, item[1], 2.18, y, 2.68, 0.24, font_size=9.5, color=COLORS["muted"], font_name="Malgun Gothic")

    add_image_or_placeholder(slide, "customer_scent_journey", ASSETS["customer_scent_journey"], 5.75, 2.08, 6.78, 4.3)
    add_footer(slide)


def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["surface"])

    add_section_title(slide, "Product Lineup", "세 가지 무드로 구성된 시그니처 컬렉션", dark=False)
    add_text(
        slide,
        "Citrus Clean, Floral Muse, Woody Noir 세 가지 향 라인은 다양한 취향과 사용 장면을 커버합니다.",
        0.74,
        2.03,
        5.1,
        0.54,
        font_size=14,
        color=COLORS["muted"],
        font_name="Malgun Gothic",
        line_spacing=1.08
    )

    add_image_or_placeholder(slide, "product_lineup_comparison", ASSETS["product_lineup_comparison"], 0.8, 2.72, 11.7, 4.1)


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])
    add_section_title(slide, "Scent Architecture", "Top, Middle, Base Note의 입체적 설계", dark=True)

    add_image_or_placeholder(slide, "scent_note_pyramid", ASSETS["scent_note_pyramid"], 0.98, 2.0, 7.15, 4.75)

    add_card(slide, 8.55, 2.18, 3.95, 3.95, fill_color="132033", line_color=COLORS["accent"])
    add_text(slide, "Fragrance Design", 8.9, 2.56, 2.6, 0.28, font_size=11, color=COLORS["accent"], font_name=FONTS["body"], bold=True)
    add_text(
        slide,
        "첫 향, 중심 향, 잔향이 시간에 따라 자연스럽게 전개되는 구조입니다.",
        8.9,
        3.08,
        3.05,
        0.65,
        font_size=15,
        color=COLORS["surface"],
        font_name="Malgun Gothic",
        bold=True,
        line_spacing=1.08
    )

    points = [
        "Top Note: 첫인상을 결정하는 산뜻한 개방감",
        "Middle Note: 브랜드 무드를 형성하는 중심 향",
        "Base Note: 오래 남는 깊이와 잔향",
    ]
    for i, point in enumerate(points):
        y = 4.1 + i * 0.45
        add_text(slide, "—", 8.9, y, 0.25, 0.24, font_size=12, color=COLORS["accent"], font_name=FONTS["body"])
        add_text(slide, point, 9.18, y, 2.88, 0.26, font_size=10.5, color="D9D0C3", font_name="Malgun Gothic")

    add_footer(slide)


def build_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["surface"])

    add_section_title(slide, "Brand Differentiation", "향, 패키지, 경험을 하나로 설계", dark=False)

    add_text(
        slide,
        "Scent Atelier는 합리적 프리미엄 가격대와 감각적 패키징, 스토리 중심 경험으로 차별화합니다.",
        0.76,
        2.25,
        4.25,
        0.86,
        font_size=18,
        color=COLORS["text"],
        font_name="Malgun Gothic",
        bold=True,
        line_spacing=1.08
    )

    diff_items = [
        ("01", "Accessible Premium", "고급감은 유지하되 진입 장벽은 낮춘 가격 전략"),
        ("02", "Sensory Packaging", "선물성과 소장 가치를 높이는 패키지 경험"),
        ("03", "Story-led Experience", "향조와 감정 장면을 연결하는 브랜드 콘텐츠"),
    ]
    for i, item in enumerate(diff_items):
        y = 3.55 + i * 0.72
        add_text(slide, item[0], 0.78, y, 0.42, 0.25, font_size=10, color=COLORS["accent"], font_name=FONTS["body"], bold=True)
        add_text(slide, item[1], 1.28, y, 2.05, 0.25, font_size=11.5, color=COLORS["background"], font_name=FONTS["body"], bold=True)
        add_text(slide, item[2], 1.28, y + 0.28, 3.55, 0.25, font_size=9.5, color=COLORS["muted"], font_name="Malgun Gothic")

    add_image_or_placeholder(slide, "brand_positioning_map", ASSETS["brand_positioning_map"], 5.35, 1.88, 6.95, 4.85)


def build_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])
    add_section_title(slide, "Next Step", "브랜드 론칭을 위한 실행 방향", dark=True)

    add_text(
        slide,
        "브랜드 론칭은 제품 검증, 캠페인 설계, 유통 확장의 3단계로 추진합니다.",
        1.05,
        2.1,
        11.1,
        0.42,
        font_size=18,
        color=COLORS["surface"],
        font_name="Malgun Gothic",
        bold=True,
        align=PP_ALIGN.CENTER
    )

    add_image_or_placeholder(slide, "launch_roadmap", ASSETS["launch_roadmap"], 1.0, 2.95, 11.33, 2.65)

    add_divider(slide, 4.85, 6.15, 3.6, color=COLORS["accent"], width=1)
    add_text(
        slide,
        "Partnership · Retail · Campaign Collaboration",
        3.35,
        6.42,
        6.65,
        0.3,
        font_size=13,
        color=COLORS["accent"],
        font_name=FONTS["body"],
        bold=True,
        align=PP_ALIGN.CENTER
    )
    add_text(
        slide,
        "contact@scentatelier.example",
        4.42,
        6.78,
        4.5,
        0.25,
        font_size=10,
        color="D9D0C3",
        font_name=FONTS["body"],
        align=PP_ALIGN.CENTER
    )
    add_corner_mark(slide)


def main():
    prs = Presentation()
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    main()