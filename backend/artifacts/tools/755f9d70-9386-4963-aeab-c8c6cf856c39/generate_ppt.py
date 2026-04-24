import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


OUTPUT_FILE = "premium_perfume_brand_intro.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "Aptos Display"
FONT_BODY = "Aptos"

COLORS = {
    "background": RGBColor(22, 18, 17),
    "surface": RGBColor(36, 28, 26),
    "accent": RGBColor(200, 164, 93),
    "secondary": RGBColor(126, 79, 58),
    "text": RGBColor(247, 241, 232),
    "muted": RGBColor(183, 169, 154),
    "black": RGBColor(0, 0, 0),
    "white": RGBColor(255, 255, 255),
}

ASSETS = {
    "hero_bottle": os.path.join("assets", "perfume_hero_bottle.png"),
    "collection_lineup": os.path.join("assets", "perfume_collection_lineup.png"),
    "scent_pyramid": os.path.join("assets", "fragrance_scent_pyramid.png"),
    "positioning_chart": os.path.join("assets", "perfume_positioning_chart.png"),
    "brand_journey": os.path.join("assets", "brand_experience_journey.png"),
}


def set_slide_background(slide, color=COLORS["background"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size=18,
    color=COLORS["text"],
    font_name=FONT_BODY,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.1,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_multi_text(
    slide,
    lines,
    x,
    y,
    w,
    h,
    font_size=18,
    color=COLORS["text"],
    font_name=FONT_BODY,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)

    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_shape(
    slide,
    shape_type,
    x,
    y,
    w,
    h,
    fill_color=None,
    line_color=None,
    line_width=1,
    transparency=0,
):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.fill.transparency = transparency
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1, y1, x2, y2, color=COLORS["accent"], width=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_card(
    slide,
    x,
    y,
    w,
    h,
    fill_color=COLORS["surface"],
    line_color=COLORS["secondary"],
    radius=True,
    transparency=0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    return add_shape(
        slide,
        shape_type,
        x,
        y,
        w,
        h,
        fill_color=fill_color,
        line_color=line_color,
        line_width=1,
        transparency=transparency,
    )


def add_image_or_placeholder(slide, asset_id, x, y, w, h, label=None):
    path = ASSETS.get(asset_id)
    if path and os.path.exists(path):
        return slide.shapes.add_picture(path, x, y, width=w, height=h)

    placeholder = add_card(
        slide,
        x,
        y,
        w,
        h,
        fill_color=COLORS["surface"],
        line_color=COLORS["accent"],
        radius=True,
        transparency=8,
    )
    add_text_box(
        slide,
        label or asset_id,
        x + Inches(0.15),
        y + h / 2 - Inches(0.22),
        w - Inches(0.3),
        Inches(0.44),
        font_size=14,
        color=COLORS["muted"],
        font_name=FONT_BODY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return placeholder


def add_slide_header(slide, title, subtitle=None):
    add_text_box(
        slide,
        title,
        Inches(0.72),
        Inches(0.42),
        Inches(5.2),
        Inches(0.45),
        font_size=22,
        color=COLORS["text"],
        font_name=FONT_TITLE,
        bold=True,
    )
    if subtitle:
        add_text_box(
            slide,
            subtitle,
            Inches(0.72),
            Inches(0.88),
            Inches(5.2),
            Inches(0.28),
            font_size=10,
            color=COLORS["accent"],
            font_name=FONT_BODY,
        )
    add_line(slide, Inches(0.72), Inches(1.26), Inches(2.1), Inches(1.26), COLORS["accent"], 1.2)


def add_footer_mark(slide, number):
    add_text_box(
        slide,
        f"{number:02d}",
        Inches(12.15),
        Inches(6.96),
        Inches(0.45),
        Inches(0.22),
        font_size=9,
        color=COLORS["muted"],
        font_name=FONT_BODY,
        align=PP_ALIGN.RIGHT,
    )


def add_decorative_circles(slide):
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        Inches(11.4),
        Inches(-0.45),
        Inches(2.6),
        Inches(2.6),
        fill_color=COLORS["secondary"],
        line_color=None,
        transparency=70,
    )
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        Inches(-0.8),
        Inches(6.25),
        Inches(1.9),
        Inches(1.9),
        fill_color=COLORS["accent"],
        line_color=None,
        transparency=84,
    )


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_decorative_circles(slide)

    add_text_box(
        slide,
        "ÉCLAT NOIR",
        Inches(0.75),
        Inches(1.42),
        Inches(6.7),
        Inches(0.82),
        font_size=48,
        color=COLORS["text"],
        font_name=FONT_TITLE,
        bold=True,
    )
    add_line(slide, Inches(0.78), Inches(2.42), Inches(2.75), Inches(2.42), COLORS["accent"], 2.2)
    add_text_box(
        slide,
        "A Premium Fragrance Brand Introduction",
        Inches(0.78),
        Inches(2.62),
        Inches(5.9),
        Inches(0.34),
        font_size=14,
        color=COLORS["accent"],
        font_name=FONT_BODY,
    )
    add_text_box(
        slide,
        "기억에 남는 순간을 향으로 설계하는 니치 퍼퓸 브랜드",
        Inches(0.78),
        Inches(3.42),
        Inches(5.9),
        Inches(0.72),
        font_size=22,
        color=COLORS["text"],
        font_name=FONT_BODY,
        bold=True,
    )
    add_text_box(
        slide,
        "ÉCLAT NOIR는 섬세한 원료와 현대적 감성을 결합해 개인의 분위기를 완성하는 프리미엄 향수 브랜드입니다.",
        Inches(0.78),
        Inches(4.32),
        Inches(5.25),
        Inches(0.9),
        font_size=14,
        color=COLORS["muted"],
        font_name=FONT_BODY,
    )
    add_image_or_placeholder(
        slide,
        "hero_bottle",
        Inches(7.35),
        Inches(0.92),
        Inches(4.65),
        Inches(5.9),
        label="hero_bottle",
    )
    add_footer_mark(slide, 1)


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_slide_header(slide, "Brand Essence", "향으로 표현하는 정체성")

    add_text_box(
        slide,
        "Deep,\nQuiet,\nMemorable",
        Inches(0.78),
        Inches(1.78),
        Inches(4.65),
        Inches(2.1),
        font_size=38,
        color=COLORS["text"],
        font_name=FONT_TITLE,
        bold=True,
    )
    add_text_box(
        slide,
        "브랜드 핵심 가치는 절제된 우아함, 오래 지속되는 잔향, 개인의 취향을 존중하는 독창성입니다.",
        Inches(0.82),
        Inches(4.16),
        Inches(4.62),
        Inches(0.9),
        font_size=15,
        color=COLORS["muted"],
        font_name=FONT_BODY,
    )

    cards = [
        ("01", "Restrained Elegance", "과장보다 절제, 유행보다 지속성을 우선하는 미감"),
        ("02", "Lasting Trail", "시간이 지나도 은은하게 남는 깊은 잔향"),
        ("03", "Personal Signature", "개인의 취향과 분위기를 완성하는 독창성"),
    ]
    for i, (num, title, body) in enumerate(cards):
        x = Inches(6.05)
        y = Inches(1.55 + i * 1.55)
        add_card(slide, x, y, Inches(5.75), Inches(1.22))
        add_text_box(slide, num, x + Inches(0.24), y + Inches(0.22), Inches(0.52), Inches(0.28), 11, COLORS["accent"], FONT_BODY, True)
        add_text_box(slide, title, x + Inches(0.92), y + Inches(0.2), Inches(3.9), Inches(0.3), 16, COLORS["text"], FONT_BODY, True)
        add_text_box(slide, body, x + Inches(0.92), y + Inches(0.62), Inches(4.35), Inches(0.38), 11, COLORS["muted"], FONT_BODY)

    add_footer_mark(slide, 2)


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_slide_header(slide, "Signature Collection", "대표 향수 라인업")

    add_text_box(
        slide,
        "세 가지 시그니처 향으로 완성하는 브랜드 세계관",
        Inches(0.75),
        Inches(1.48),
        Inches(9.2),
        Inches(0.48),
        font_size=20,
        color=COLORS["text"],
        font_name=FONT_BODY,
        bold=True,
    )
    add_image_or_placeholder(slide, "collection_lineup", Inches(0.82), Inches(2.02), Inches(11.65), Inches(1.15), label="collection_lineup")

    products = [
        ("Noir Wood", "Woody Musk", "TOP  Bergamot\nHEART  Cedarwood\nBASE  Musk, Amber", "DEEP · WARM"),
        ("Lumière Rose", "Floral Amber", "TOP  Pink Pepper\nHEART  Rose, Iris\nBASE  Amber, Vanilla", "SOFT · LUMINOUS"),
        ("Citrus Veil", "Citrus Musk", "TOP  Mandarin\nHEART  Neroli\nBASE  White Musk", "FRESH · CLEAN"),
    ]

    for i, (name, family, notes, mood) in enumerate(products):
        x = Inches(0.82 + i * 4.13)
        y = Inches(3.45)
        add_card(slide, x, y, Inches(3.62), Inches(2.55))
        add_shape(slide, MSO_SHAPE.OVAL, x + Inches(1.44), y + Inches(0.2), Inches(0.72), Inches(0.72), COLORS["secondary"], COLORS["accent"], 1, 15)
        add_text_box(slide, name, x + Inches(0.25), y + Inches(1.02), Inches(3.12), Inches(0.33), 17, COLORS["text"], FONT_BODY, True, align=PP_ALIGN.CENTER)
        add_text_box(slide, family, x + Inches(0.25), y + Inches(1.42), Inches(3.12), Inches(0.25), 10, COLORS["accent"], FONT_BODY, align=PP_ALIGN.CENTER)
        add_multi_text(slide, notes.split("\n"), x + Inches(0.42), y + Inches(1.8), Inches(2.8), Inches(0.55), 9, COLORS["muted"], FONT_BODY, False, PP_ALIGN.CENTER)
        add_text_box(slide, mood, x + Inches(0.25), y + Inches(2.28), Inches(3.12), Inches(0.22), 8, COLORS["accent"], FONT_BODY, True, align=PP_ALIGN.CENTER)

    add_footer_mark(slide, 3)


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_slide_header(slide, "Scent Architecture", "향의 구조")

    add_text_box(
        slide,
        "첫 향부터 잔향까지 설계된 감각의 흐름",
        Inches(0.76),
        Inches(1.5),
        Inches(7.8),
        Inches(0.45),
        font_size=20,
        color=COLORS["text"],
        font_name=FONT_BODY,
        bold=True,
    )
    add_text_box(
        slide,
        "상큼한 탑 노트, 풍성한 하트 노트, 깊고 지속적인 베이스 노트가 균형 있게 이어집니다.",
        Inches(0.78),
        Inches(2.02),
        Inches(4.95),
        Inches(0.65),
        font_size=13,
        color=COLORS["muted"],
        font_name=FONT_BODY,
    )

    add_image_or_placeholder(slide, "scent_pyramid", Inches(4.58), Inches(1.82), Inches(4.2), Inches(4.15), label="scent_pyramid")

    labels = [
        ("Top Note", "Mandarin · Bergamot · Pink Pepper", Inches(8.95), Inches(2.03)),
        ("Heart Note", "Rose · Neroli · Cedarwood · Iris", Inches(9.34), Inches(3.45)),
        ("Base Note", "Musk · Amber · Vanilla · Woods", Inches(8.92), Inches(4.92)),
    ]
    for title, body, x, y in labels:
        add_card(slide, x, y, Inches(3.38), Inches(0.78), fill_color=COLORS["surface"], line_color=COLORS["accent"])
        add_text_box(slide, title, x + Inches(0.2), y + Inches(0.12), Inches(2.8), Inches(0.22), 12, COLORS["text"], FONT_BODY, True)
        add_text_box(slide, body, x + Inches(0.2), y + Inches(0.42), Inches(2.9), Inches(0.2), 8.5, COLORS["muted"], FONT_BODY)
        add_line(slide, x - Inches(0.42), y + Inches(0.39), x, y + Inches(0.39), COLORS["accent"], 1)

    add_footer_mark(slide, 4)


def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_slide_header(slide, "Target Audience", "핵심 고객층")

    add_text_box(
        slide,
        "취향을 통해 자신을 표현하는 25–40세 프리미엄 소비자",
        Inches(0.78),
        Inches(1.5),
        Inches(5.3),
        Inches(0.75),
        font_size=20,
        color=COLORS["text"],
        font_name=FONT_BODY,
        bold=True,
    )
    add_text_box(
        slide,
        "핵심 고객은 대중적인 향보다 개성 있는 니치 향을 선호하며, 패키지 디자인과 브랜드 스토리까지 구매 기준으로 고려합니다.",
        Inches(0.8),
        Inches(2.42),
        Inches(4.95),
        Inches(0.85),
        font_size=13,
        color=COLORS["muted"],
        font_name=FONT_BODY,
    )

    add_card(slide, Inches(0.82), Inches(3.62), Inches(4.9), Inches(1.72))
    add_text_box(slide, "Persona", Inches(1.05), Inches(3.88), Inches(1.9), Inches(0.28), 13, COLORS["accent"], FONT_BODY, True)
    add_multi_text(
        slide,
        ["도시적 라이프스타일", "감각적 소비와 자기표현", "리미티드 제품과 브랜드 스토리 선호"],
        Inches(1.05),
        Inches(4.28),
        Inches(4.25),
        Inches(0.78),
        font_size=11,
        color=COLORS["muted"],
        font_name=FONT_BODY,
    )

    add_image_or_placeholder(slide, "positioning_chart", Inches(6.55), Inches(1.52), Inches(5.65), Inches(4.75), label="positioning_chart")

    add_footer_mark(slide, 5)


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_slide_header(slide, "Brand Experience", "고객 접점 전략")

    add_text_box(
        slide,
        "향을 발견하고, 경험하고, 기억하게 만드는 여정",
        Inches(0.78),
        Inches(1.45),
        Inches(8.2),
        Inches(0.55),
        font_size=21,
        color=COLORS["text"],
        font_name=FONT_BODY,
        bold=True,
    )
    add_text_box(
        slide,
        "팝업 부티크, 시향 키트, 맞춤 추천 콘텐츠, 멤버십 리워드를 통해 고객 경험을 확장합니다.",
        Inches(0.8),
        Inches(2.05),
        Inches(7.45),
        Inches(0.42),
        font_size=13,
        color=COLORS["muted"],
        font_name=FONT_BODY,
    )

    steps = [
        ("Discovery", "SNS 콘텐츠\n시향 키트"),
        ("Experience", "팝업 부티크\n컨설팅"),
        ("Conversion", "프리미엄 패키지\n온라인 구매"),
        ("Loyalty", "멤버십 리워드\n리필 프로그램"),
    ]
    for i, (title, body) in enumerate(steps):
        x = Inches(0.88 + i * 3.05)
        y = Inches(3.05)
        add_card(slide, x, y, Inches(2.4), Inches(1.45))
        add_text_box(slide, f"0{i + 1}", x + Inches(0.18), y + Inches(0.18), Inches(0.42), Inches(0.22), 10, COLORS["accent"], FONT_BODY, True)
        add_text_box(slide, title, x + Inches(0.18), y + Inches(0.48), Inches(1.95), Inches(0.28), 14, COLORS["text"], FONT_BODY, True)
        add_text_box(slide, body, x + Inches(0.18), y + Inches(0.88), Inches(1.95), Inches(0.42), 10, COLORS["muted"], FONT_BODY)
        if i < len(steps) - 1:
            add_line(slide, x + Inches(2.42), y + Inches(0.72), x + Inches(3.0), y + Inches(0.72), COLORS["accent"], 1.4)

    add_image_or_placeholder(slide, "brand_journey", Inches(0.88), Inches(5.08), Inches(11.6), Inches(1.12), label="brand_journey")
    add_footer_mark(slide, 6)


def build_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_slide_header(slide, "Growth Strategy", "시장 확장 방향")

    add_text_box(
        slide,
        "니치 향수 시장에서 감각적 브랜드 자산을 축적하는 단계적 성장",
        Inches(0.78),
        Inches(1.48),
        Inches(9.7),
        Inches(0.65),
        font_size=21,
        color=COLORS["text"],
        font_name=FONT_BODY,
        bold=True,
    )
    add_text_box(
        slide,
        "1단계 브랜드 론칭, 2단계 리테일 및 팝업 확장, 3단계 글로벌 온라인 채널 진출을 추진합니다.",
        Inches(0.8),
        Inches(2.17),
        Inches(7.4),
        Inches(0.45),
        font_size=13,
        color=COLORS["muted"],
        font_name=FONT_BODY,
    )

    phases = [
        ("01", "Launch", "브랜드 아이덴티티 확립\n시그니처 3종 출시\nD2C 온라인몰 오픈"),
        ("02", "Expand", "프리미엄 편집숍 입점\n팝업 부티크 운영\n시향 키트 캠페인"),
        ("03", "Scale", "글로벌 온라인 채널 진출\n한정판 컬렉션 확장\n멤버십 데이터 고도화"),
    ]
    y = Inches(3.25)
    for i, (num, title, body) in enumerate(phases):
        x = Inches(0.88 + i * 4.05)
        add_card(slide, x, y, Inches(3.5), Inches(2.08))
        add_text_box(slide, num, x + Inches(0.28), y + Inches(0.26), Inches(0.55), Inches(0.3), 12, COLORS["accent"], FONT_BODY, True)
        add_text_box(slide, title, x + Inches(0.28), y + Inches(0.68), Inches(2.5), Inches(0.38), 18, COLORS["text"], FONT_TITLE, True)
        add_multi_text(slide, body.split("\n"), x + Inches(0.28), y + Inches(1.2), Inches(2.82), Inches(0.68), 10.5, COLORS["muted"], FONT_BODY)
        if i < len(phases) - 1:
            add_line(slide, x + Inches(3.52), y + Inches(1.04), x + Inches(4.0), y + Inches(1.04), COLORS["accent"], 1.4)

    add_text_box(
        slide,
        "ÉCLAT NOIR builds fragrance as a lasting brand asset.",
        Inches(2.4),
        Inches(6.18),
        Inches(8.5),
        Inches(0.35),
        font_size=16,
        color=COLORS["accent"],
        font_name=FONT_BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer_mark(slide, 7)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    main()