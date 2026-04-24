from __future__ import annotations

from textwrap import dedent

from app.models.schemas import PPTCodeBundle, PlanningSpec, VisualAssetsBundle


def generate_ppt_code(planning_spec: PlanningSpec, visual_assets_bundle: VisualAssetsBundle) -> PPTCodeBundle:
    # LoRA3 입력 계약의 핵심: 시각 에셋 레이아웃 코드 포함
    embedded_assets = [
        {
            "asset_id": a.asset_id,
            "slide_index": a.slide_index,
            "layout_code": a.layout_code,
            "placement": a.placement,
        }
        for a in visual_assets_bundle.assets
    ]

    code = dedent(
        f"""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        VISUAL_ASSETS_LAYOUT_CODE = {embedded_assets!r}

        prs = Presentation()
        for idx, slide_spec in enumerate({[s.model_dump() for s in planning_spec.slides]!r}, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.5), Inches(0.8))
            tf = title_box.text_frame
            tf.text = slide_spec["title"]
            tf.paragraphs[0].font.size = Pt(30)

            asset = next((a for a in VISUAL_ASSETS_LAYOUT_CODE if a["slide_index"] == idx), None)
            if asset:
                p = asset["placement"]
                box = slide.shapes.add_textbox(Inches(p["x"]), Inches(p["y"]), Inches(p["w"]), Inches(p["h"]))
                box.text_frame.text = f'Asset: {{asset["asset_id"]}}'

        prs.save("output.pptx")
        """
    ).strip()

    return PPTCodeBundle(python_code=code, expected_outputs=["output.pptx"])


def repair_ppt_code(prior_code: str, traceback_text: str) -> PPTCodeBundle:
    repaired = prior_code + f"\n# Repair hint: {traceback_text[:160]}\n"
    return PPTCodeBundle(python_code=repaired, expected_outputs=["output.pptx"])
